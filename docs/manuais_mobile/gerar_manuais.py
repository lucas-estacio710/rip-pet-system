# -*- coding: utf-8 -*-
# Gera os 2 manuais mobile (PPTX) — Fluxo PV e Recontratação de Tutor.
# Prints em shots/; saída na raiz do projeto.
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SHOTS = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'shots')
OUT_DIR = r'C:\Users\kel_v\LUCAS_ATELLIAR\NOVO_CRM_RIP_PET'

AZUL = RGBColor(0x1E, 0x3A, 0x5F)
AZUL_CLARO = RGBColor(0x2A, 0x78, 0xD6)
VERDE = RGBColor(0x19, 0x9E, 0x70)
CINZA = RGBColor(0x47, 0x55, 0x69)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def novo_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def slide_capa(prs, titulo, subtitulo):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fundo = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    fundo.fill.solid(); fundo.fill.fore_color.rgb = AZUL; fundo.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.6))
    p = tb.text_frame.paragraphs[0]; p.text = titulo
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = BRANCO
    tb2 = s.shapes.add_textbox(Inches(0.9), Inches(4.1), Inches(11.5), Inches(1.2))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = subtitulo
    p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(0xB8, 0xD4, 0xF0)
    tb3 = s.shapes.add_textbox(Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5))
    p3 = tb3.text_frame.paragraphs[0]; p3.text = 'R.I.P. Pet — CRM · Manual do usuário (mobile)'
    p3.font.size = Pt(13); p3.font.color.rgb = RGBColor(0x8F, 0xA8, 0xC4)

def slide_texto(prs, titulo, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.9), Inches(0.6), Inches(11.5), Inches(1.0))
    p = tb.text_frame.paragraphs[0]; p.text = titulo
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = AZUL
    tb2 = s.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(5.2))
    tf = tb2.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.text = b
        par.font.size = Pt(17); par.font.color.rgb = CINZA
        par.space_after = Pt(12)
    return s

def slide_passo(prs, numero, titulo, img, bullets, destaque=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Print mobile à esquerda (razão 390x844)
    img_path = os.path.join(SHOTS, img + '.png')
    alt = Inches(6.9)
    larg = Emu(int(alt * 390 / 844))
    s.shapes.add_picture(img_path, Inches(0.7), Inches(0.3), height=alt)
    # Número do passo
    x_txt = Inches(4.6)
    tb0 = s.shapes.add_textbox(x_txt, Inches(0.55), Inches(1.6), Inches(1.0))
    p0 = tb0.text_frame.paragraphs[0]; p0.text = f'{numero:02d}'
    p0.font.size = Pt(48); p0.font.bold = True; p0.font.color.rgb = AZUL_CLARO
    # Título
    tb = s.shapes.add_textbox(x_txt + Inches(1.5), Inches(0.72), Inches(6.7), Inches(1.4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = titulo
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = AZUL
    # Bullets
    tb2 = s.shapes.add_textbox(x_txt, Inches(2.3), Inches(8.0), Inches(4.2))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for i, b in enumerate(bullets):
        par = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        par.text = '•  ' + b
        par.font.size = Pt(16); par.font.color.rgb = CINZA
        par.space_after = Pt(10)
    # Destaque (caixa verde)
    if destaque:
        box = s.shapes.add_shape(1, x_txt, Inches(6.1), Inches(8.0), Inches(1.0))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE6, 0xF7, 0xF0)
        box.line.color.rgb = VERDE
        tfd = box.text_frame; tfd.word_wrap = True
        pd = tfd.paragraphs[0]; pd.text = '💡 ' + destaque
        pd.font.size = Pt(14); pd.font.bold = True; pd.font.color.rgb = RGBColor(0x0F, 0x6B, 0x4A)
    return s

# ================= MANUAL 1 — FLUXO PV =================
prs = novo_deck()
slide_capa(prs, 'Contratação e Ativação de Plano Preventivo (PV)',
           'Do link público que o tutor preenche até o acionamento do plano')

slide_texto(prs, 'Visão geral do fluxo', [
    '1. O tutor recebe o link público da unidade (/preventivo/…) e preenche a contratação com o pet AINDA VIVO.',
    '2. A ficha chega em Fichas de Entrada com a faixa lateral VERDE (PV) — diferente da vermelha (emergencial).',
    '3. O concierge processa a ficha: valida o contato, confere o valor do plano e cria o contrato.',
    '4. O contrato nasce com status PREVENTIVO e fica na tela Preventivos (não conta como atendimento ativo).',
    '5. Quando o pet falecer, o operador clica em ATV (Acionar PV): preenche o acolhimento e o contrato vira ATIVO, seguindo o fluxo normal de remoção/cremação.',
])

slide_passo(prs, 1, 'Tutor abre o link público da unidade', 'pv-01-form-tutor', [
    'Cada unidade tem seu próprio link: rippet.com.br/preventivo/santos, /sao-paulo, /campinas, etc.',
    'O formulário é público — o tutor preenche do próprio celular, sem login.',
    'Passos do formulário: 1 Tutor → 2 Pet → 3 Confirmar.',
])

slide_passo(prs, 2, 'Passo 1 — Dados do tutor', 'pv-02-form-tutor-preenchido', [
    'Nome completo (vai pro contrato e certificado), CPF/CNPJ, telefone e e-mail.',
    'CEP com busca automática de endereço ("Buscar Endereço").',
    'O CPF é validado na hora — número inválido bloqueia o envio.',
    'O rascunho fica salvo no aparelho ("Rascunho salvo") — se o tutor sair, não perde o que digitou.',
])

slide_passo(prs, 3, 'Passo 2 — Dados do pet (vivo!)', 'pv-03-form-pet', [
    'Nome, idade, espécie (Canina/Felina/Exótica), gênero, raça, cor e peso aproximado.',
    'Escolha do plano: Individual (COM retorno das cinzas) ou Coletiva (SEM retorno).',
    'Não há campos de óbito: localização, velório e acompanhamento não aparecem no formulário preventivo.',
])

slide_passo(prs, 4, 'Passo 3 — Confirmação', 'pv-04-confirmacao', [
    'Resumo de tudo que foi preenchido: dados do tutor e do pet.',
    'O tutor confere e pode voltar pra corrigir qualquer coisa.',
])

slide_passo(prs, 5, '"Como conheceu" e envio', 'pv-05-como-conheceu', [
    'O tutor marca como conheceu a R.I.P. Pet (Google, Instagram, indicação…).',
    'Tutores antigos costumam marcar "Já utilizei a R.I.P. Pet" — isso aparece depois pro concierge.',
    'Botão verde "Enviar Ficha ✓" conclui.',
])

slide_passo(prs, 6, 'Ficha enviada', 'pv-06-sucesso', [
    'Tela de sucesso confirma o registro.',
    'O tutor é avisado que receberá confirmação por WhatsApp.',
    'A partir daqui, o fluxo segue DENTRO do CRM.',
])

slide_passo(prs, 7, 'Ficha chega em Fichas de Entrada', 'pv-07-fichas-recebida', [
    'O card entra como "Recebida" com a FAIXA LATERAL VERDE escrita PV — contratação preventiva.',
    'Fichas de remoção (emergencial) têm faixa vermelha EM.',
    'Toque em "Processar" pra abrir a tratativa.',
], destaque='Faixa verde = pet vivo. Não é remoção!'),

slide_passo(prs, 8, 'Processar Ficha — tipo PREVENTIVO', 'pv-08-tratativa-tutor-existente', [
    'O cabeçalho mostra "DADOS DA FICHA — TIPO: PREVENTIVO".',
    'Se o CPF já existe na base, aparece o selo verde "Tutor existente" — o contrato será vinculado ao cadastro antigo, sem duplicar.',
    'Confira os dados do tutor e do pet (lápis permite corrigir).',
])

slide_passo(prs, 9, 'Valide o contato para cremação', 'pv-09-contato-cremacao', [
    'A Matriz usará esse número pra chamar no dia da cremação.',
    'Confirme: "Sim, é este" (e diga como chamar o contato) ou "Não, é outro" (informe o número certo).',
    'É obrigatório validar o contato pra processar a ficha.',
])

slide_passo(prs, 10, 'Valor do plano e desconto', 'pv-10-valor-plano', [
    'Toque num valor sugerido (990, 1.090, …, 1.690) ou digite o valor fechado.',
    'Desconto pré-venda (R$ ou %) se houver parceria/condição especial.',
    '"Detalhamento do Plano" vira a descrição no PDF do contrato (ex: Plano Gratidão, com molde da patinha e urna MDF inclusas).',
])

slide_passo(prs, 11, 'Ficha processada', 'pv-11-ficha-processada', [
    'Depois do "Processar", o card fica verde ("Processada") e mostra o valor do plano.',
    'A ficha ainda NÃO virou contrato — falta o último passo.',
])

slide_passo(prs, 12, 'Criar o contrato', 'pv-12-criar-contrato', [
    'Abra "Visualizar Ficha" e confira o resumo final.',
    '"Enviar Confirmação" manda a mensagem de boas-vindas por WhatsApp.',
    '"Gerar Contrato PDF" emite o contrato preventivo pra assinatura/aceite.',
    'Toque em "Criar Contrato" pra efetivar.',
], destaque='O contrato PV usa um PDF próprio — diferente do emergencial.')

slide_passo(prs, 13, 'Contrato nasce como PREVENTIVO', 'pv-13-pipeline-preventivo', [
    'O contrato aparece no pipeline com o selo dourado "Preventivo" e as tags PV + IND/COL.',
    'Código gerado automaticamente (ex: ST260719INDTUTREXUY).',
    'PV não conta como atendimento ativo — o pet está vivo.',
])

slide_passo(prs, 14, 'Tela Preventivos — gestão dos planos', 'pv-14-preventivos-card', [
    'Todos os contratos preventivos da unidade, com busca por pet, tutor ou código.',
    'Farol de pagamento: verde = pago, amarelo = parcial, vermelho = a pagar.',
    'Botão WhatsApp fala direto com o tutor.',
    'O botão ATV aciona o plano quando o pet falece.',
])

slide_passo(prs, 15, 'Acionar PV — o pet faleceu', 'pv-15-ativar-modal', [
    'O modal "Ativar Contrato" abre o mesmo Acolhimento da tratativa emergencial.',
    'Valide o contato para cremação (o mesmo do plano ou outro).',
])

slide_passo(prs, 16, 'Preencha o acolhimento e ative', 'pv-16-ativar-acolhimento', [
    'Local do acolhimento (Unidade R.I.P. PET ou outro endereço), responsável, data/hora e número do lacre.',
    'Ao confirmar, o contrato passa de PREVENTIVO pra ATIVO e entra no fluxo normal (encaminhamento → cremação → retorno).',
], destaque='Depois de acionado, o PV não volta a ser preventivo.')

prs.save(os.path.join(OUT_DIR, 'Manual_Fluxo_PV.pptx'))
print('Manual_Fluxo_PV.pptx ok —', len(prs.slides.__iter__.__self__._sldIdLst), 'slides')

# ================= MANUAL 2 — RECONTRATAÇÃO =================
prs2 = novo_deck()
slide_capa(prs2, 'Recontratação de Tutor Antigo',
           'Novo contrato pra quem já é cliente — sem duplicar cadastro')

slide_texto(prs2, 'Visão geral do fluxo', [
    '1. O tutor que já usou a R.I.P. Pet quer contratar de novo (novo pet, emergência ou plano preventivo).',
    '2. Busque o cadastro dele em Tutores e use "Enviar nova contratação" — o link já vai com os dados preenchidos.',
    '3. O tutor só confirma os dados e cadastra o NOVO pet.',
    '4. Na tratativa, o CRM reconhece o CPF e mostra "Tutor existente" — o novo contrato entra no MESMO cadastro.',
    '5. O histórico do tutor acumula todos os contratos (antigos e novos).',
])

slide_passo(prs2, 1, 'Busque o tutor', 'rt-01-busca-tutor', [
    'Tela Tutores: busque por nome, telefone, CPF, e-mail ou cidade.',
    'Confirme que é a mesma pessoa pelo CPF/telefone antes de seguir.',
])

slide_passo(prs2, 2, 'Abra o cadastro do tutor', 'rt-02-cadastro-tutor', [
    'O cadastro mostra contatos, endereço e o total de contratos.',
    'Toque em "Enviar nova contratação".',
], destaque='Esse botão é o caminho certo pra recontratação — evita cadastro duplicado.')

slide_passo(prs2, 3, 'Envie o link pré-preenchido', 'rt-03-link-nova-contratacao', [
    'O CRM gera um link exclusivo com os dados de cadastro do tutor já preenchidos.',
    'Ele só confirma e cadastra o novo pet. O link vale 7 dias.',
    'Envie por WhatsApp com um toque, ou copie o link.',
])

slide_passo(prs2, 4, 'A ficha volta com "Tutor existente"', 'pv-08-tratativa-tutor-existente', [
    'Quando a nova ficha chega e você processa, o CRM reconhece o CPF.',
    'O selo verde "Tutor existente" confirma: o contrato será vinculado ao cadastro antigo.',
    'NUNCA crie outro cadastro pro mesmo CPF — o vínculo é automático.',
])

slide_passo(prs2, 5, 'Histórico acumula os contratos', 'rt-04-historico-contratos', [
    'No cadastro do tutor, o Histórico de Contratos lista todos os pets e contratos, antigos e novos.',
    'Endereço com atalho pra Waze e Google Maps.',
])

slide_texto(prs2, 'Pontos de atenção', [
    '✅ Sempre busque o tutor ANTES de mandar link novo — se ele já existe, use "Enviar nova contratação".',
    '✅ O reconhecimento é pelo CPF: se a ficha vier com CPF diferente (digitado errado), o CRM cria OUTRO cadastro. Corrija o CPF na tratativa (lápis) antes de processar.',
    '✅ Tutores que marcam "Já utilizei a R.I.P. Pet" no formulário aparecem sinalizados na tratativa.',
    '✅ O telefone/endereço novos informados na ficha atualizam o contato do contrato — o cadastro do tutor pode ser editado na tela Tutores.',
])

prs2.save(os.path.join(OUT_DIR, 'Manual_Recontratacao_Tutor.pptx'))
print('Manual_Recontratacao_Tutor.pptx ok')
